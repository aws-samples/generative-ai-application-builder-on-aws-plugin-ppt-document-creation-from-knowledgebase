import os
import re
from typing import List, Optional

import pptx
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from bisheng.utils import defaults
from bisheng.utils.defaults import (
    SYSTEM_PROMPT_PATTERN,
    DEFAULT_SYSTEM_PROMPT,
    GENERATE_PROMPT_PATTERN,
    OUTPUT_PROMPT_PATTERN
)

from bisheng.decoders.base_decoder import BaseDecoder
from bisheng.prompting.base_prompt import BasePrompt
from bisheng.prompting.templates.template_factory import TemplateFactory


def _has_system_prompt(slide) -> bool:
    return len(re.findall(SYSTEM_PROMPT_PATTERN, slide.notes_slide.notes_text_frame.text.strip(), re.DOTALL)) > 0


def _get_system_prompt(slide) -> Optional[str]:
    return re.findall(SYSTEM_PROMPT_PATTERN, slide.notes_slide.notes_text_frame.text.strip(), re.DOTALL)[0] \
        if _has_system_prompt(slide) \
        else DEFAULT_SYSTEM_PROMPT


def _has_generate_prompt(shape) -> bool:
    return len(re.findall(GENERATE_PROMPT_PATTERN, shape.text_frame.text.strip(), re.DOTALL)) > 0


def _get_generate_prompt(shape) -> Optional[str]:
    return re.findall(GENERATE_PROMPT_PATTERN, shape.text_frame.text.strip(), re.DOTALL)[0]


def _has_output_prompt(shape) -> bool:
    return len(re.findall(OUTPUT_PROMPT_PATTERN, shape.text_frame.text.strip(), re.DOTALL)) > 0


def _get_output_prompt(shape) -> Optional[str]:
    return re.findall(OUTPUT_PROMPT_PATTERN, shape.text_frame.text.strip(), re.DOTALL)[0]


class PptxDecoder(BaseDecoder):
    class PptxPromptBehavior(BasePrompt):

        def __init__(self, slide: Slide, shape: BaseShape, system_prompt, instruction):
            self.slide = slide
            self.shape = shape
            self.system_prompt = system_prompt
            self.instruction = instruction

        def to_json(self):
            return {
                "instruction": self.instruction,
                "system_prompt": self.system_prompt,
                "slide_id": self.slide.slide_id,
                "slide_name": self.slide.name,
                "shape_id": self.shape.shape_id,
                "shape_name": self.shape.name
            }

        def get_system_prompt(self):
            return self.system_prompt

        def get_instruction_prompt(self):
            return self.instruction

        def get_slide(self):
            return self.slide

        def get_shape(self):
            return self.shape

    def __init__(self,
                 prompts_path: str,
                 shots_path: str,
                 context_path: str,
                 instruction: str
                 ):
        self.prompts_path = prompts_path
        self.shots_path = shots_path
        self.context_path = context_path
        self.instruction = instruction
        self._process_files(context_path, prompts_path, shots_path)

        self.prompts: List[PptxDecoder.PptxPromptBehavior] = []
        self.template_factory = TemplateFactory(config=instruction)

    def _create_prompt(self, shape, slide, instruction_template) -> PptxPromptBehavior:
        system_prompt = _get_system_prompt(slide)
        generate_prompt = _get_generate_prompt(shape)
        output_filter = _get_output_prompt(shape) if _has_output_prompt(shape) else defaults.DEFAULT_OUTPUT_INDICATOR

        prompt_slides = self.prompts_pptx.slides
        shot_slides = self.shots_pptx.slides
        shot_slide_idx = prompt_slides.index(slide)
        shot_slide = shot_slides[shot_slide_idx]
        shot_shape_idx = slide.shapes.index(shape)
        shot_shape = shot_slide.shapes[shot_shape_idx]
        shot_text = shot_shape.text_frame.text.strip()

        all_params = {
            "generate_prompt": generate_prompt,
            "shot_text": shot_text,
            "context": self.context,
            "output_filter": output_filter
        }
        required_params = instruction_template.get_required_params()
        instruction_parameters = {k: v for k, v in all_params.items() if k in required_params}
        instruction = instruction_template.create_instruction(**instruction_parameters)
        return PptxDecoder.PptxPromptBehavior(slide, shape, system_prompt, instruction)

    def _find_prompts(self):
        slides = self.prompts_pptx.slides
        for slide in slides:
            qualified_shapes = list(filter(lambda shape: shape.has_text_frame, slide.shapes))
            qualified_shapes = list(filter(_has_generate_prompt, qualified_shapes))
            self.prompts.extend([self._create_prompt(shape, slide, self.template_factory.create()) for shape in
                                 qualified_shapes])
        self.num_prompts = len(self.prompts)

    def _process_files(self, context_path, prompts_path, shots_path):
        if os.path.isfile(prompts_path):
            self.prompts_pptx = pptx.Presentation(prompts_path)
        else:
            raise FileNotFoundError(f"File not found: {prompts_path}")
        if os.path.isfile(shots_path):
            self.shots_pptx = pptx.Presentation(shots_path)
        else:
            raise FileNotFoundError(f"File not found: {shots_path}")
        if os.path.isfile(context_path):
            with open(context_path, 'r') as file:
                self.context = file.read()
        else:
            raise FileNotFoundError(f"File not found: {context_path}")

    def decode(self, *args, **input_data: dict):
        self._find_prompts()

    def get_encoder_metadata(self) -> dict:
        encoder_metadata = {
            "context": self.context,
            "slides": self.prompts_pptx.slides
        }
        return encoder_metadata
