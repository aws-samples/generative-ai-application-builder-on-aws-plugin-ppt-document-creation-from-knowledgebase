import json
import os

import pptx
from pptx.util import Pt

from bisheng.encoders.base_encoder import BaseEncoder


class PptxEncoder(BaseEncoder):

    def __init__(self, path: str, append: str):
        self.path = path
        self.append = append

        if os.path.isfile(self.path):
            self.artifact_pptx = pptx.Presentation(self.path)
        else:
            raise FileNotFoundError(f"File not found: {self.path}")

    def encode(self, **input_data):
        results = input_data["results"]
        slides = input_data["slides"]

        artifact_slides = self.artifact_pptx.slides
        for prompt, response in results.items():
            artifact_slide_idx = slides.index(prompt.get_slide())
            artifact_slide = artifact_slides[artifact_slide_idx]
            artifact_shape_idx = prompt.get_slide().shapes.index(prompt.get_shape())
            artifact_shape = artifact_slide.shapes[artifact_shape_idx]

            artifact_slide.notes_slide.notes_text_frame.text += response.rationale
            artifact_slide.notes_slide.notes_text_frame.text += "\n\n\n"
            artifact_slide.notes_slide.notes_text_frame.text += response.sources

            if self.append:
                paragraph = artifact_shape.text_frame.add_paragraph()
                paragraph.font.size = Pt(12)
                paragraph.text = response.response.strip()
            else:
                artifact_shape.text_frame.text = response.response.strip()

        self.artifact_pptx.save(self.path)
